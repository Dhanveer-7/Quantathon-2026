"""
═══════════════════════════════════════════════════════════════════════════════
  RIT Quantathan 2026 — PowerPoint Generator (Expanded 12-Slide Pitch)
  File: populate_ppt.py
  Team: Schrödinger’s Coders
═══════════════════════════════════════════════════════════════════════════════
"""

import os
import sys
# Configure stdout to support UTF-8 characters in Windows console
if sys.stdout.encoding != 'utf-8':
    try:
        sys.stdout.reconfigure(encoding='utf-8')
    except Exception:
        pass

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PPT_PATH = "RIT - Quantathan 2026 - PPT.pptx"

def main():
    if not os.path.exists(PPT_PATH):
        print(f"[ERROR] PPT template file not found at: {PPT_PATH}")
        return
        
    print("[PPT] Loading PowerPoint template...")
    prs = Presentation(PPT_PATH)
    
    # ─── SLIDE 0: TITLE SLIDE ───
    slide_0 = prs.slides[0]
    
    # Modify text of existing shapes to preserve positioning
    for shape in slide_0.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text
        if "Presented by" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = "Q-GUARD AI"
            p.font.name = 'Rajdhani'
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 212, 255) # Glowing Cyan
            
            p2 = shape.text_frame.add_paragraph()
            p2.text = "Multi-Agent Quantum Fraud Intelligence System"
            p2.font.name = 'Exo 2'
            p2.font.size = Pt(18)
            p2.font.color.rgb = RGBColor(220, 240, 255) # Light Ice Blue
            p2.space_before = Pt(8)
            
        elif "Name of the Candidate" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = "Team Schrödinger’s Coders"
            p.font.name = 'Exo 2'
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 255, 157) # Bright Green
            
            p2 = shape.text_frame.add_paragraph()
            p2.text = "Members: Dhinith Pragalyan, Dhanveer M, Venkatesan S, Elamaran B"
            p2.font.name = 'Exo 2'
            p2.font.size = Pt(12)
            p2.font.color.rgb = RGBColor(240, 240, 240) # White
            p2.space_before = Pt(6)
            
        elif "Institution / Organization" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = "RIT Chennai — Quantathan 2026"
            p.font.name = 'Share Tech Mono'
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(200, 200, 200)

    # ─── CLEAR ALL SLIDES EXCEPT TITLE SLIDE 0 ───
    id_list = prs.slides._sldIdLst
    while len(prs.slides) > 1:
        del id_list[1]

    # ─── EXPANDED 12 CONTENT SLIDES DATA ───
    content_layout = prs.slide_layouts[1]
    
    slides_data = [
        {
            "title": "1. PROBLEM STATEMENT (QT-2.8)",
            "bullets": [
                "Mule Fund Relays: Syndicates route stolen funds through intermediate accounts in minutes.",
                "Device Farms: Fraudsters use virtual emulators to spoof 'trusted devices' at scale.",
                "Velocity Structuring: Bots slice transactions into micro-amounts below default detection limits.",
                "Linear Blindspot: Traditional rules evaluate details in isolation, missing non-linear loops.",
                "Solution: Q-Guard AI blends Multi-Agent ingestion with a Parameterized Quantum Circuit."
            ]
        },
        {
            "title": "2. WHY CLASSICAL MACHINE LEARNING FAILS",
            "bullets": [
                "High False Positives: Static rules trigger blocks on regular traveler activity, creating churn.",
                "Dimensionality Curse: Linking deep behavior details requires heavy classical networks.",
                "Black-Box Decisioning: Deep learning models cannot explain logs for audit compliance.",
                "Rigid Threshold Limits: Static fraud limits leave banks open to off-peak night exploits."
            ]
        },
        {
            "title": "3. THE QUANTUM ADVANTAGE (QML LAYER)",
            "bullets": [
                "Non-Linear Correlation: Quantum state entanglements capture inter-agent correlations naturally.",
                "Parameter Efficiency: A 4-qubit classifier uses 90% fewer weights than classical models.",
                "Hilbert Mapping: Translates complex user behavior metrics to subatomic quantum states.",
                "Hardware Readiness: Designed to plug natively into physical NISQ QPUs (like IBM Quantum)."
            ]
        },
        {
            "title": "4. MULTI-AGENT INGESTION ARCHITECTURE",
            "bullets": [
                "User Agent: Analyzes spending limit fraction deviations and transaction timestamps.",
                "Device Agent: Assesses browser attributes, OS versions, emulator footprints, and device age.",
                "Merchant Agent: Monitors seller category hazard levels and historic chargeback metrics.",
                "Network Agent: Assesses VPN/proxy usage, geodistance anomalies, and impossible travel velocity.",
                "Decentralized Design: Ensures multiple threat dimensions are ingested simultaneously."
            ]
        },
        {
            "title": "5. BLOCH SPHERE & STATE VECTOR MAPPING",
            "bullets": [
                "Angle Encoding: Local agent threat outputs (0 to 1) map to rotation angles (0 to Pi).",
                "Qubit Allocation: Initializes a 4-qubit state, applying Ry rotation gates to each qubit.",
                "Bloch Projections: Risk levels are mapped as coordinates on a 3D Bloch Sphere vector.",
                "Excitation Tracking: Elevated telemetry flags push the state vectors closer to the |1> pole."
            ]
        },
        {
            "title": "6. VARIATIONAL QUANTUM CIRCUIT (VQC) DETAILS",
            "bullets": [
                "Ansatz Entanglement: Connects qubits in a ring topology (Q0->Q1->Q2->Q3->Q0) using CNOT gates.",
                "Coordinated Learning: Entanglement allows Q-Guard to associate user details with device risk.",
                "Variational Tuning: Fine-tunes model weights dynamically using rotation parameter loops.",
                "Output Expectation: Calculates the Z-axis expectation value of each qubit upon measurement."
            ]
        },
        {
            "title": "7. DYNAMIC BLOCH VECTOR WEIGHTING",
            "bullets": [
                "Dynamic Weights: Calculates weights directly from Bloch Z-projection: W = (1 - Z)/2.",
                "Adaptive Focus: System dynamically shifts priority to the agent showing the highest anomaly.",
                "Combined Risk Scoring: Calculates final score as a weighted sum of qubit measurements.",
                "Noise Suppression: Low-risk background telemetry is dynamically ignored to save database bandwidth."
            ]
        },
        {
            "title": "8. ADAPTIVE CONTEXT THRESHOLD ENGINE",
            "bullets": [
                "Dynamic Guardrails: Replaces flat, static 80% limits with context-aware boundaries.",
                "Time Window (12 AM - 5 AM): Threshold drops to 65% to defend against off-hour scams.",
                "Merchant Context: Limit drops to 58% for high-hazard category terminal devices.",
                "Customer Context: Threshold relaxes to 95% for trusted, long-standing user profiles."
            ]
        },
        {
            "title": "9. SELF-EXPLAINABLE AI (XAI) & AUDITING",
            "bullets": [
                "Translation Pipeline: Converts VQC coordinates into natural language audit logs.",
                "Regulatory Compliance: Outlines exactly why a transaction was flagged or blocked.",
                "Anti-Black Box: Ensures fraud analysts can explain any system decision to RBI auditors.",
                "Visual Threat Mapping: Provides a dynamic 3D visual of the Bloch sphere state space."
            ]
        },
        {
            "title": "10. STRATEGIC MITIGATION CHECKLISTS",
            "bullets": [
                "Multi-Stage Security: Replaces crude binary block decisions with progressive checklists.",
                "Low Risk: Process transaction instantly (frictionless flow).",
                "Medium Risk: Prompt for SMS OTP and email code verification.",
                "High Risk: Require biometric Face Authentication check.",
                "Critical Risk: Freeze transaction and alert the bank's fraud mitigation desk."
            ]
        },
        {
            "title": "11. SCALABILITY & DATABASE INTEGRATION",
            "bullets": [
                "Dataset Loading: Mapped to 3,000 real-world transaction logs from the Kaggle dataset.",
                "MongoDB Pipeline: Stores telemetry, dynamic weights, and strategies in indexed collections.",
                "FastAPI REST Endpoints: Exposes Bloch coordinates, agent weights, and explanation logs.",
                "Vercel Deployment: Live static visualizer deployed on secure cloud servers."
            ]
        },
        {
            "title": "12. FUTURE SCOPE & CONCLUSION",
            "bullets": [
                "QPU Deployment: Porting the simulated quantum layers directly to physical IBM QPUs.",
                "Graph Neural Network (GNN): Merging VQC with GNNs to trace large money laundering rings.",
                "Open Banking API: Packaging Q-Guard as a plug-and-play middleware module.",
                "Schrödinger's Coders: Securing the future of payments using quantum computing principles."
            ]
        }
    ]

    for sd in slides_data:
        slide = prs.slides.add_slide(content_layout)
        
        # Set Title
        slide.shapes.title.text = sd["title"]
        
        # Populate Content Placeholder
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.word_wrap = True
        tf.clear()
        
        for idx, bullet_text in enumerate(sd["bullets"]):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = bullet_text
            p.level = 0
            p.space_after = Pt(14)
            # Inherit template fonts, colors, and bullet markers automatically for maximum readability

    prs.save(PPT_PATH)
    print(f"[PPT] ✅ Successfully compiled PowerPoint presentation with {len(prs.slides)} slides!")

if __name__ == "__main__":
    main()
